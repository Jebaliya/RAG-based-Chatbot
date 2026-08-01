"""
ingestion.py
============
Stage 1 of the pipeline: reads raw uploaded files (PDF, DOCX, PPTX, TXT)
and loads each into LangChain `Document` objects -- one per page/slide/
section, with `source_file` and `page_number` set in `.metadata` so later
stages can cite exactly where an answer came from.

USES LANGCHAIN DOCUMENT LOADERS:
- PDFPlumberLoader (langchain_community) -- wraps pdfplumber, one Document
  per PDF page. Chosen over PyPDFLoader because the project already
  standardized on pdfplumber's extraction quality.
- Docx2txtLoader (langchain_community) -- wraps docx2txt, a small pure-
  Python DOCX text extractor with no heavy transitive dependencies.
- TextLoader (langchain_community) -- plain .txt files.
- PPTXLoader (below) -- a small CUSTOM loader for PowerPoint files, built
  directly on python-pptx. LangChain's own PPTX loader
  (UnstructuredPowerPointLoader) pulls in the large `unstructured` package
  and often needs system-level dependencies (e.g. libreoffice) -- too
  heavy for a Streamlit Community Cloud deployment. Implementing
  LangChain's `BaseLoader` interface directly keeps the same PPTX-parsing
  logic the project already had, while still slotting into the standard
  LangChain loader pattern (`.load()` -> `list[Document]`).
"""

from pathlib import Path

from langchain_community.document_loaders import PDFPlumberLoader, Docx2txtLoader, TextLoader
from langchain_core.document_loaders import BaseLoader
from langchain_core.documents import Document
from pptx import Presentation


class PPTXLoader(BaseLoader):
    """
    Custom LangChain-compatible loader for PowerPoint files. Extracts each
    slide's text into its own Document, tagged with `source_file` and
    `page_number` (the slide number) metadata -- kept consistent with the
    other loaders below.
    """

    def __init__(self, filepath: Path):
        self.filepath = Path(filepath)

    def load(self) -> list[Document]:
        documents = []
        prs = Presentation(self.filepath)
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in paragraph.runs)
                        if line.strip():
                            texts.append(line)
            slide_text = "\n".join(texts)
            if slide_text.strip():
                documents.append(Document(
                    page_content=slide_text,
                    metadata={"source_file": self.filepath.name, "page_number": i},
                ))
        return documents


def _load_pdf(filepath: Path) -> list[Document]:
    docs = PDFPlumberLoader(str(filepath)).load()
    for i, doc in enumerate(docs, start=1):
        doc.metadata["source_file"] = filepath.name
        doc.metadata["page_number"] = doc.metadata.get("page", i - 1) + 1
    return [d for d in docs if d.page_content.strip()]


def _load_docx(filepath: Path) -> list[Document]:
    docs = Docx2txtLoader(str(filepath)).load()
    for doc in docs:
        doc.metadata["source_file"] = filepath.name
        doc.metadata["page_number"] = 1  # DOCX has no fixed page concept at the XML level
    return [d for d in docs if d.page_content.strip()]


def _load_txt(filepath: Path) -> list[Document]:
    docs = TextLoader(str(filepath), encoding="utf-8", autodetect_encoding=True).load()
    for doc in docs:
        doc.metadata["source_file"] = filepath.name
        doc.metadata["page_number"] = 1
    return [d for d in docs if d.page_content.strip()]


def _load_pptx(filepath: Path) -> list[Document]:
    return PPTXLoader(filepath).load()


_LOADERS = {
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".txt": _load_txt,
    ".pptx": _load_pptx,
}


def load_document(filepath: Path) -> list[Document]:
    """Dispatches a single file to the correct LangChain loader based on its extension."""
    loader = _LOADERS.get(filepath.suffix.lower())
    if loader is None:
        raise ValueError(f"Unsupported file type: {filepath.suffix}")
    return loader(filepath)


def load_all_documents(data_dir: Path) -> list[Document]:
    """
    Walks the upload directory and loads every supported file into a flat
    list of LangChain Documents, ready for chunking.py to split.
    """
    all_docs: list[Document] = []
    for filepath in sorted(data_dir.glob("*")):
        if filepath.suffix.lower() in _LOADERS:
            try:
                all_docs.extend(load_document(filepath))
            except Exception as exc:
                # Don't let one corrupt/unreadable file kill the whole build --
                # skip it and let the caller surface a warning to the user.
                print(f"[ingestion] Failed to load {filepath.name}: {exc}")
    return all_docs


if __name__ == "__main__":
    from config import DATA_DIR
    docs = load_all_documents(DATA_DIR)
    print(f"Loaded {len(docs)} pages/sections from {DATA_DIR}")
    for d in docs[:3]:
        print(f"--- {d.metadata['source_file']} (page {d.metadata['page_number']}) ---")
        print(d.page_content[:200], "...\n")
